# Contains the function to extract URLs from PPTX files.
import pptx
import os
from .utils import URL_PATTERN, clean_urls

def extract_urls_from_pptx(filepath):

    urls = []
    try:
        prs = pptx.Presentation(filepath)
        #Iterate through all the sldies
        for slide in prs.slides:
            #Check shapes on the slide
            for shape in slide.shapes:
                #Check for any hyperlinks on each slide
                if shape.click_action.hyperlink and shape.click_action.hyperlink.address:
                    urls.append(shape.click_action.hyperlink.address)

                #Checking for texts in shapes
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            #Hyperlink checker
                            if run.hyperlink and run.hyperlink.address:
                                urls.append(run.hyperlink.address)

            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame

    except Exception as e:
        print(f"--> Error processing PPTX file '{os.path.basename(filepath)}': {e}")
        return []
    cleaned_list = clean_urls(urls)
    return cleaned_list

